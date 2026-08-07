from pathlib import Path

import pytest
from openpyxl import Workbook

from corp_os.services.extract import SUPPORTED_EXTENSIONS, extract_text
from corp_os.services.ingest import gate_file


def test_supported_extensions_are_enterprise_allowlist():
    assert SUPPORTED_EXTENSIONS == {
        ".txt",
        ".md",
        ".csv",
        ".xlsx",
        ".pdf",
        ".png",
        ".jpg",
        ".jpeg",
        ".webp",
        ".docx",
        ".pptx",
    }
    # Legacy / high-risk formats stay out.
    for ext in {".doc", ".ppt", ".xls", ".xlsm", ".mp4", ".mov", ".zip", ".exe", ".gif", ".bmp"}:
        assert ext not in SUPPORTED_EXTENSIONS


def test_extract_xlsx(tmp_path: Path):
    path = tmp_path / "demo.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "费用"
    ws.append(["项目", "金额"])
    ws.append(["差旅", 1200])
    wb.save(path)

    text = extract_text(path)
    assert "费用" in text
    assert "差旅" in text
    assert "1200" in text


def test_extract_unsupported_video_returns_empty(tmp_path: Path):
    path = tmp_path / "demo.mp4"
    path.write_bytes(b"\x00\x00fake-video")
    assert extract_text(path) == ""


def test_extract_override(tmp_path: Path):
    path = tmp_path / "a.png"
    path.write_bytes(b"not-a-real-image")
    assert extract_text(path, override="  手填发票金额 100  ") == "手填发票金额 100"


def test_extract_pptx(tmp_path: Path):
    from pptx import Presentation

    path = tmp_path / "template.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "项目汇报模板"
    slide.placeholders[1].text = "背景\n方案"
    prs.save(str(path))

    text = extract_text(path)
    assert "幻灯片" in text
    assert "项目汇报模板" in text
    assert "背景" in text


def test_gate_rejects_unsupported_extension(tmp_path: Path):
    class _FakeUpload:
        filename = "clip.mp4"
        content_type = "video/mp4"
        file = None

        def __init__(self, data: bytes):
            from io import BytesIO

            self.file = BytesIO(data)

    with pytest.raises(ValueError, match="不支持的文件类型"):
        gate_file(_FakeUpload(b"fake"))
