from pathlib import Path

from corp_os.config import get_settings
from corp_os.services.docgen import (
    generate_markdown,
    generate_powerpoint,
    generate_word,
    resolve_generated_file,
)


def test_generate_word_and_download(tmp_path, monkeypatch):
    monkeypatch.setenv("CORP_OS_UPLOAD_DIR", str(tmp_path))
    get_settings.cache_clear()

    doc = generate_word(
        title="周报",
        sections=[
            {"heading": "本周进展", "bullets": ["完成登录", "联调 ERP"]},
            {"heading": "风险", "paragraphs": ["无重大风险"]},
        ],
        username="alice",
    )
    assert doc.path.is_file()
    assert doc.path.stat().st_size > 1000
    assert doc.download_url.endswith(doc.file_id)

    path, filename, media = resolve_generated_file(doc.file_id, username="alice")
    assert path == doc.path
    assert filename.endswith(".docx")
    assert "wordprocessingml" in media

    try:
        resolve_generated_file(doc.file_id, username="other")
        assert False, "expected PermissionError"
    except PermissionError:
        pass

    get_settings.cache_clear()


def test_generate_powerpoint(tmp_path, monkeypatch):
    monkeypatch.setenv("CORP_OS_UPLOAD_DIR", str(tmp_path))
    get_settings.cache_clear()

    doc = generate_powerpoint(
        title="项目汇报",
        slides=[
            {"title": "背景", "bullets": ["客户需求", "交付目标"]},
            {"title": "计划", "bullets": ["一期", "二期"]},
        ],
        username="boss",
    )
    assert doc.path.is_file()
    assert doc.path.suffix == ".pptx"
    assert Path(doc.path).stat().st_size > 5000

    get_settings.cache_clear()



def test_generate_powerpoint_from_template(tmp_path, monkeypatch):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    monkeypatch.setenv("CORP_OS_UPLOAD_DIR", str(tmp_path))
    get_settings.cache_clear()

    tpl = tmp_path / "tpl.pptx"
    src = Presentation()
    src.slide_width = 12192000
    src.slide_height = 6858000
    # Cover + content pages with Chinese placeholder hints (like company templates).
    cover = src.slides.add_slide(src.slide_layouts[6])  # blank
    box = cover.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    box.text_frame.text = "商务汇报标题"
    p1 = src.slides.add_slide(src.slide_layouts[6])
    t1 = p1.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    t1.text_frame.text = "添加标题"
    b1 = p1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
    b1.text_frame.text = "单击此处添加正文"
    p2 = src.slides.add_slide(src.slide_layouts[6])
    t2 = p2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    t2.text_frame.text = "添加标题"
    b2 = p2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
    b2.text_frame.text = "单击此处添加正文"
    src.save(str(tpl))

    doc = generate_powerpoint(
        title="季度汇报",
        slides=[
            {"title": "进展", "bullets": ["完成 A", "推进 B"]},
            {"title": "计划", "bullets": ["下季度重点"]},
        ],
        username="boss",
        template_path=tpl,
    )
    out = Presentation(str(doc.path))
    assert doc.path.is_file() and doc.path.stat().st_size > 2000
    assert out.slide_width == 12192000
    assert doc.template_mode == "visual"
    blob = "\n".join(
        (sh.text_frame.text or "")
        for slide in out.slides
        for sh in slide.shapes
        if getattr(sh, "has_text_frame", False)
    )
    assert "进展" in blob and "完成 A" in blob
    assert all(len(list(s.shapes)) > 0 for s in out.slides)

    get_settings.cache_clear()



def test_generate_markdown(tmp_path, monkeypatch):
    monkeypatch.setenv("CORP_OS_UPLOAD_DIR", str(tmp_path))
    get_settings.cache_clear()

    doc = generate_markdown(
        title="接口说明",
        sections=[
            {"heading": "概述", "paragraphs": ["本文说明登录接口。"]},
            {"heading": "步骤", "bullets": ["获取 token", "调用业务 API"]},
        ],
        username="alice",
    )
    text = doc.path.read_text(encoding="utf-8")
    assert doc.path.suffix == ".md"
    assert "# 接口说明" in text
    assert "## 概述" in text
    assert "- 获取 token" in text
    _, filename, media = resolve_generated_file(doc.file_id, username="alice")
    assert filename.endswith(".md")
    assert "markdown" in media

    get_settings.cache_clear()
