"""Generate Word / PowerPoint / Markdown files from structured outlines (agent tools)."""

from __future__ import annotations

import json
import re
import uuid
from dataclasses import dataclass
from pathlib import Path

from corp_os.config import get_settings

_SAFE_ID = re.compile(r"^[a-f0-9]{8,32}$", re.IGNORECASE)

_MEDIA_TYPES = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "md": "text/markdown; charset=utf-8",
}


@dataclass
class GeneratedDoc:
    file_id: str
    filename: str
    path: Path
    media_type: str
    download_url: str
    # powerpoint: "visual" = filled company template; "size" = blank deck + template size; "blank"
    template_mode: str = "blank"


def _generated_dir() -> Path:
    root = get_settings().upload_dir / "generated"
    root.mkdir(parents=True, exist_ok=True)
    return root


def _safe_stem(title: str, fallback: str) -> str:
    raw = (title or fallback).strip() or fallback
    cleaned = re.sub(r'[\\/:*?"<>|\s]+', "_", raw)[:40].strip("._")
    return cleaned or fallback


def _write_meta(file_id: str, *, username: str, filename: str, kind: str) -> None:
    meta = {
        "file_id": file_id,
        "username": username,
        "filename": filename,
        "kind": kind,
    }
    (_generated_dir() / f"{file_id}.json").write_text(
        json.dumps(meta, ensure_ascii=False),
        encoding="utf-8",
    )


def load_generated_meta(file_id: str) -> dict | None:
    if not _SAFE_ID.match(file_id or ""):
        return None
    path = _generated_dir() / f"{file_id}.json"
    if not path.is_file():
        return None
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
        return data if isinstance(data, dict) else None
    except (OSError, json.JSONDecodeError):
        return None


def resolve_generated_file(file_id: str, *, username: str, elevated: bool = False) -> tuple[Path, str, str]:
    """Return (path, filename, media_type). Raises ValueError/PermissionError."""
    meta = load_generated_meta(file_id)
    if not meta:
        raise ValueError("文件不存在或已过期")
    owner = str(meta.get("username") or "")
    if owner != username and not elevated:
        raise PermissionError("无权下载该文件")
    kind = str(meta.get("kind") or "docx")
    filename = str(meta.get("filename") or f"{file_id}.{kind}")
    path = _generated_dir() / f"{file_id}.{kind}"
    if not path.is_file():
        raise ValueError("文件不存在或已过期")
    media = _MEDIA_TYPES.get(kind, "application/octet-stream")
    return path, filename, media


def _normalize_sections(sections: list | None, body: str | None) -> list[dict]:
    out: list[dict] = []
    for item in sections or []:
        if isinstance(item, str):
            text = item.strip()
            if text:
                out.append({"heading": "", "paragraphs": [text]})
            continue
        if not isinstance(item, dict):
            continue
        heading = str(item.get("heading") or item.get("title") or "").strip()
        paragraphs: list[str] = []
        raw_paras = item.get("paragraphs") or item.get("content") or item.get("body")
        if isinstance(raw_paras, str) and raw_paras.strip():
            paragraphs.extend([p.strip() for p in re.split(r"\n+", raw_paras) if p.strip()])
        elif isinstance(raw_paras, list):
            paragraphs.extend([str(p).strip() for p in raw_paras if str(p).strip()])
        bullets = item.get("bullets") or item.get("points") or []
        if isinstance(bullets, list):
            paragraphs.extend([f"• {str(b).strip()}" for b in bullets if str(b).strip()])
        if heading or paragraphs:
            out.append({"heading": heading, "paragraphs": paragraphs})
    if not out and (body or "").strip():
        paras = [p.strip() for p in re.split(r"\n{2,}", body.strip()) if p.strip()]
        out.append({"heading": "", "paragraphs": paras or [body.strip()]})
    return out


def _sections_to_markdown(title: str, blocks: list[dict]) -> str:
    lines: list[str] = [f"# {title}", ""]
    for block in blocks:
        if block["heading"]:
            lines.append(f"## {block['heading']}")
            lines.append("")
        for para in block["paragraphs"]:
            text = para
            if text.startswith("• "):
                text = "- " + text[2:]
            lines.append(text)
            lines.append("")
    return "\n".join(lines).rstrip() + "\n"


def generate_markdown(
    *,
    title: str,
    sections: list | None = None,
    body: str | None = None,
    username: str,
) -> GeneratedDoc:
    """Write a UTF-8 Markdown (.md) file."""
    title = (title or "未命名文档").strip()[:120] or "未命名文档"
    raw_body = (body or "").strip()

    if sections:
        blocks = _normalize_sections(sections, None)
        content = _sections_to_markdown(title, blocks) if blocks else f"# {title}\n\n"
        if raw_body and raw_body not in content:
            content = content.rstrip() + "\n\n---\n\n" + raw_body + "\n"
    elif raw_body:
        if not raw_body.lstrip().startswith("#"):
            content = f"# {title}\n\n{raw_body}\n"
        else:
            content = raw_body if raw_body.endswith("\n") else raw_body + "\n"
    else:
        content = _sections_to_markdown(
            title,
            [{"heading": "", "paragraphs": ["（正文为空，请补充内容后重新生成）"]}],
        )

    file_id = uuid.uuid4().hex[:16]
    filename = f"{_safe_stem(title, 'document')}.md"
    path = _generated_dir() / f"{file_id}.md"
    path.write_text(content, encoding="utf-8")
    _write_meta(file_id, username=username, filename=filename, kind="md")
    return GeneratedDoc(
        file_id=file_id,
        filename=filename,
        path=path,
        media_type=_MEDIA_TYPES["md"],
        download_url=f"/api/v1/chat/generated/{file_id}",
    )


def generate_word(
    *,
    title: str,
    sections: list | None = None,
    body: str | None = None,
    username: str,
) -> GeneratedDoc:
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError as exc:
        raise RuntimeError("缺少 python-docx，请 pip install python-docx") from exc

    title = (title or "未命名文档").strip()[:120] or "未命名文档"
    blocks = _normalize_sections(sections, body)
    if not blocks:
        blocks = [{"heading": "", "paragraphs": ["（正文为空，请补充内容后重新生成）"]}]

    doc = Document()
    doc.add_heading(title, level=0)
    for block in blocks:
        if block["heading"]:
            doc.add_heading(block["heading"], level=1)
        for para in block["paragraphs"]:
            p = doc.add_paragraph(para)
            for run in p.runs:
                run.font.size = Pt(11)

    file_id = uuid.uuid4().hex[:16]
    filename = f"{_safe_stem(title, 'document')}.docx"
    path = _generated_dir() / f"{file_id}.docx"
    doc.save(str(path))
    _write_meta(file_id, username=username, filename=filename, kind="docx")
    return GeneratedDoc(
        file_id=file_id,
        filename=filename,
        path=path,
        media_type=_MEDIA_TYPES["docx"],
        download_url=f"/api/v1/chat/generated/{file_id}",
    )


def _pptx_output_ok(path: Path, *, title: str, slides: list[dict]) -> bool:
    """Reject corrupt / empty-shape decks (commercial template deletion failure mode)."""
    try:
        from pptx import Presentation

        if not path.is_file() or path.stat().st_size < 2000:
            return False
        prs = Presentation(str(path))
        if len(prs.slides) == 0:
            return False
        for slide in prs.slides:
            if len(list(slide.shapes)) == 0:
                return False
        blob = "\n".join(
            _shape_plain_text(sh)
            for slide in prs.slides
            for sh in _iter_text_shapes(slide)
        )
        markers = [title] + [str(s.get("title") or "") for s in slides[:3]]
        markers = [m for m in markers if m and m.strip()]
        return any(m in blob for m in markers) if markers else True
    except Exception:  # noqa: BLE001
        return False


def generate_powerpoint(
    *,
    title: str,
    slides: list | None = None,
    sections: list | None = None,
    username: str,
    template_path: Path | str | None = None,
) -> GeneratedDoc:
    """Prefer filling the company template in-place; fall back to a stable fresh deck.

    Never deletes slides from commercial templates (that previously corrupted OOXML).
    """
    import logging
    import shutil

    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("缺少 python-pptx，请 pip install python-pptx") from exc

    log = logging.getLogger(__name__)
    title = (title or "未命名演示").strip()[:120] or "未命名演示"
    raw_slides = slides if slides is not None else sections
    normalized: list[dict] = []
    for item in raw_slides or []:
        if isinstance(item, str):
            text_item = item.strip()
            if text_item:
                normalized.append({"title": text_item[:40], "bullets": [text_item]})
            continue
        if not isinstance(item, dict):
            continue
        slide_title = str(item.get("title") or item.get("heading") or "").strip() or "要点"
        bullets: list[str] = []
        for key in ("bullets", "points", "paragraphs"):
            val = item.get(key)
            if isinstance(val, list):
                bullets.extend([str(x).strip() for x in val if str(x).strip()])
            elif isinstance(val, str) and val.strip():
                bullets.extend([p.strip() for p in re.split(r"\n+", val) if p.strip()])
        content = item.get("content")
        if isinstance(content, str) and content.strip():
            bullets.extend([p.strip() for p in re.split(r"\n+", content) if p.strip()])
        if not bullets:
            bullets = ["（请补充内容）"]
        normalized.append({"title": slide_title[:80], "bullets": bullets[:12]})

    if not normalized:
        normalized = [{"title": title, "bullets": ["（幻灯片内容为空，请补充后重新生成）"]}]

    file_id = uuid.uuid4().hex[:16]
    filename = f"{_safe_stem(title, 'presentation')}.pptx"
    out_dir = _generated_dir()
    path = out_dir / f"{file_id}.pptx"

    template_mode = "blank"
    tpl = Path(template_path) if template_path else None
    used_visual = False

    if tpl is not None and tpl.is_file():
        try:
            shutil.copy2(tpl, path)
            prs = Presentation(str(path))
            if _apply_template_safely(prs, deck_title=title, slides=normalized):
                prs.save(str(path))
                if _pptx_output_ok(path, title=title, slides=normalized):
                    used_visual = True
                    template_mode = "visual"
                else:
                    log.warning("template fill produced invalid pptx; falling back to fresh deck")
            else:
                log.info("template had no fillable placeholders; falling back to fresh deck")
        except Exception:  # noqa: BLE001
            log.exception("template fill failed; falling back to fresh deck")
            used_visual = False

    if not used_visual:
        prs = Presentation()
        if tpl is not None and tpl.is_file():
            try:
                src = Presentation(str(tpl))
                prs.slide_width = src.slide_width
                prs.slide_height = src.slide_height
                template_mode = "size"
            except Exception:  # noqa: BLE001
                template_mode = "blank"
        _build_deck_from_layouts(prs, deck_title=title, slides=normalized)
        prs.save(str(path))

    if not path.is_file() or path.stat().st_size < 2000:
        raise RuntimeError("PPT 写入失败：生成文件不存在或过小")

    _write_meta(file_id, username=username, filename=filename, kind="pptx")
    if load_generated_meta(file_id) is None:
        raise RuntimeError("PPT 元数据写入失败")

    return GeneratedDoc(
        file_id=file_id,
        filename=filename,
        path=path,
        media_type=_MEDIA_TYPES["pptx"],
        download_url=f"/api/v1/chat/generated/{file_id}",
        template_mode=template_mode,
    )


_SAFE_CJK_FONTS = ("微软雅黑", "Microsoft YaHei", "PingFang SC", "华文黑体", "SimHei")

_PLACEHOLDER_HINTS = (
    "添加标题",
    "单击此处",
    "标题内容",
    "预祝工作",
    "可以直接复制",
    "内容要与标题",
    "内容需要和标题",
    "内容要复合",
    "内容要与标题一致",
    "点击此处",
    "在此输入",
    "Click to edit",
    "click to edit",
    "单击此处编辑",
)


def _ensure_cjk_font(run, typeface: str | None = None) -> None:
    try:
        from lxml import etree
        from pptx.oxml.ns import qn
    except ImportError:
        return
    name = (typeface or "").strip() or _SAFE_CJK_FONTS[0]
    try:
        run.font.name = name
    except Exception:  # noqa: BLE001
        pass
    try:
        rPr = run._r.get_or_add_rPr()
        rPr.set("lang", "zh-CN")
        rPr.set("altLang", "en-US")
        for tag in ("a:latin", "a:ea", "a:cs"):
            el = rPr.find(qn(tag))
            if el is None:
                el = etree.SubElement(rPr, qn(tag))
            el.set("typeface", name)
            if "charset" not in el.attrib:
                el.set("charset", "-122")
    except Exception:  # noqa: BLE001
        pass


def _iter_text_shapes(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            yield shape


def _shape_plain_text(shape) -> str:
    try:
        return shape.text_frame.text or ""
    except Exception:  # noqa: BLE001
        return ""


def _is_placeholder_text(text: str) -> bool:
    t = (text or "").strip()
    if not t:
        return False
    return any(h in t for h in _PLACEHOLDER_HINTS)


def _is_decorative_label(text: str) -> bool:
    t = (text or "").strip()
    if not t:
        return True
    if len(t) <= 3 and t.replace(" ", "").isdigit():
        return True
    low = t.lower()
    if low in {"business", "work report", "01", "02", "03", "04", "0 1", "0 2", "0 3", "0 4"}:
        return True
    if "ppt模板" in low or "www." in low or "http" in low:
        return True
    if any(k in t for k in ("素材", "授权", "禁止商用", "谢谢观看", "WORK REPORT", "ENTERPRISE")):
        return True
    return False


def _clone_run_font_to(target_run, source_run) -> None:
    try:
        from copy import deepcopy

        from pptx.oxml.ns import qn

        src_rPr = source_run._r.find(qn("a:rPr"))
        if src_rPr is None:
            return
        tgt_r = target_run._r
        old = tgt_r.find(qn("a:rPr"))
        if old is not None:
            tgt_r.remove(old)
        tgt_r.insert(0, deepcopy(src_rPr))
    except Exception:  # noqa: BLE001
        pass


def _first_styled_run(text_frame):
    from pptx.oxml.ns import qn

    for p in text_frame.paragraphs:
        for r in p.runs:
            if (r.text or "").strip() or r._r.find(qn("a:rPr")) is not None:
                return r
        if p.runs:
            return p.runs[0]
    return None


def _replace_shape_text(shape, lines: list[str]) -> None:
    """Replace text in-place; never clear rPr. Only touch this shape."""
    text_frame = shape.text_frame
    lines = [str(x) for x in (lines or []) if str(x).strip()] or [""]
    style_run = _first_styled_run(text_frame)
    paragraphs = list(text_frame.paragraphs)
    while len(paragraphs) < len(lines):
        text_frame.add_paragraph()
        paragraphs = list(text_frame.paragraphs)

    for i, para in enumerate(paragraphs):
        if i >= len(lines):
            for r in para.runs:
                r.text = ""
            continue
        text = lines[i]
        runs = list(para.runs)
        if not runs:
            run = para.add_run()
            if style_run is not None:
                _clone_run_font_to(run, style_run)
            run.text = text
            _ensure_cjk_font(run, _SAFE_CJK_FONTS[0])
        else:
            if style_run is not None:
                _clone_run_font_to(runs[0], style_run)
            runs[0].text = text
            _ensure_cjk_font(runs[0], _SAFE_CJK_FONTS[0])
            for r in runs[1:]:
                r.text = ""


def _delete_slide(prs, index: int) -> None:
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001
    sld_id = sld_id_lst[index]
    prs.part.drop_rel(sld_id.rId)
    del sld_id_lst[index]


def _delete_all_slides(prs) -> None:
    while len(prs.slides) > 0:
        _delete_slide(prs, len(prs.slides) - 1)


def _write_placeholder_body(shape, lines: list[str]) -> None:
    from pptx.util import Pt

    tf = shape.text_frame
    # Prefer writing into existing paragraphs/runs to keep style.
    _replace_shape_text(shape, lines)
    for p in tf.paragraphs:
        for r in p.runs:
            if r.font.size is None:
                try:
                    r.font.size = Pt(16)
                except Exception:  # noqa: BLE001
                    pass
            _ensure_cjk_font(r, _SAFE_CJK_FONTS[0])


def _fill_only_placeholders(slide, *, title: str, bullets: list[str]) -> bool:
    """Return True if at least one placeholder was updated. Never touch decorative text.

    Unused placeholder boxes are left unchanged (do not blank them — blanking + slide
    deletion previously corrupted commercial templates into background-only pages).
    """
    title_boxes: list = []
    body_boxes: list = []
    for shape in _iter_text_shapes(slide):
        raw = _shape_plain_text(shape).strip()
        is_ph = False
        ph_idx = None
        try:
            is_ph = bool(shape.is_placeholder)
            if is_ph:
                ph_idx = int(shape.placeholder_format.idx)
        except Exception:  # noqa: BLE001
            is_ph = False
        fillable = _is_placeholder_text(raw) or (
            is_ph and (not raw or len(raw) < 80) and not _is_decorative_label(raw)
        )
        if not fillable:
            continue
        # idx 0/title-ish → title; others → body
        if raw == "添加标题" or (raw.startswith("添加标题") and len(raw) <= 12) or ph_idx == 0:
            title_boxes.append(shape)
        else:
            body_boxes.append(shape)

    if not title_boxes and not body_boxes:
        return False

    clean = [b for b in bullets if str(b).strip()]
    changed = False

    if title and title_boxes:
        _replace_shape_text(title_boxes[0], [title])
        changed = True
    elif title and body_boxes:
        _replace_shape_text(body_boxes[0], [title])
        changed = True
        body_boxes = body_boxes[1:]

    if body_boxes:
        if len(body_boxes) == 1:
            if clean:
                _write_placeholder_body(body_boxes[0], clean)
                changed = True
        else:
            for i, shape in enumerate(body_boxes):
                if i < len(clean):
                    _write_placeholder_body(shape, [clean[i]])
                    changed = True
                # else: keep original placeholder text

    return changed


def _set_cover_main_title(slide, deck_title: str, subtitle: str | None = None) -> bool:
    """Change only the main Chinese title on cover; leave English decorations alone."""
    candidates = []
    for shape in _iter_text_shapes(slide):
        raw = _shape_plain_text(shape).strip()
        if not raw or _is_decorative_label(raw):
            continue
        if _is_placeholder_text(raw):
            candidates.append((0, shape, raw))
            continue
        if any("\u4e00" <= ch <= "\u9fff" for ch in raw) and 2 <= len(raw) <= 40:
            area = int(getattr(shape, "width", 0) or 0) * int(getattr(shape, "height", 0) or 0)
            candidates.append((area, shape, raw))
    if not candidates:
        return False
    candidates.sort(key=lambda x: x[0], reverse=True)
    _replace_shape_text(candidates[0][1], [deck_title])
    return True


def _content_slide_indices(prs) -> list[int]:
    idxs: list[int] = []
    for i, slide in enumerate(prs.slides):
        blob = "\n".join(_shape_plain_text(sh) for sh in _iter_text_shapes(slide))
        if any(k in blob for k in ("素材授权", "禁止商用", "可以在下列情况使用", "PPT模板：")):
            continue
        if "PPT模板 http" in blob or blob.strip().startswith("PPT模板"):
            continue
        if _is_decorative_label(blob) and not any(h in blob for h in _PLACEHOLDER_HINTS):
            continue
        if any(h in blob for h in _PLACEHOLDER_HINTS):
            idxs.append(i)
    return idxs


def _thanks_slide_index(prs) -> int | None:
    for i, slide in enumerate(prs.slides):
        blob = "\n".join(_shape_plain_text(sh) for sh in _iter_text_shapes(slide))
        if "谢谢观看" in blob:
            return i
    return None


def _build_deck_from_layouts(prs, *, deck_title: str, slides: list[dict]) -> None:
    """Create clean title + content slides on an empty presentation."""
    from pptx.util import Pt

    title_layout = prs.slide_layouts[0]
    body_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else title_layout

    cover = prs.slides.add_slide(title_layout)
    if cover.shapes.title is not None:
        cover.shapes.title.text = deck_title
        for r in cover.shapes.title.text_frame.paragraphs[0].runs:
            _ensure_cjk_font(r, _SAFE_CJK_FONTS[0])
    if len(cover.placeholders) > 1 and slides:
        sub = " · ".join(slides[0].get("bullets") or [])[:80] or "由 corp-os 生成"
        cover.placeholders[1].text = sub
        for r in cover.placeholders[1].text_frame.paragraphs[0].runs:
            _ensure_cjk_font(r, _SAFE_CJK_FONTS[0])

    for item in slides:
        slide = prs.slides.add_slide(body_layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = item["title"]
            for r in slide.shapes.title.text_frame.paragraphs[0].runs:
                _ensure_cjk_font(r, _SAFE_CJK_FONTS[0])
        body = None
        try:
            body = slide.placeholders[1]
        except Exception:  # noqa: BLE001
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape != getattr(slide.shapes, "title", None):
                    body = shape
                    break
        if body is None:
            continue
        tf = body.text_frame
        tf.clear()
        bullets = item.get("bullets") or [" "]
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = bullet
            p.level = 0
            run.font.size = Pt(18)
            _ensure_cjk_font(run, _SAFE_CJK_FONTS[0])


def _apply_template_safely(prs, *, deck_title: str, slides: list[dict]) -> bool:
    """Fill company template in-place. Never delete slides (deletion corrupts many .pptx).

    Returns True if at least one content slide was filled.
    """
    content_idxs = _content_slide_indices(prs)
    if not content_idxs:
        return False

    if len(prs.slides) > 0 and deck_title:
        _set_cover_main_title(prs.slides[0], deck_title)

    fill_targets = [i for i in content_idxs if i != 0] or list(content_idxs)
    filled = 0
    for i, item in enumerate(slides):
        if i >= len(fill_targets):
            break
        ok = _fill_only_placeholders(
            prs.slides[fill_targets[i]],
            title=item.get("title") or "",
            bullets=item.get("bullets") or [],
        )
        if ok:
            filled += 1
    return filled > 0


def format_tool_result(doc: GeneratedDoc, kind_label: str) -> str:
    return (
        f"已生成{kind_label}《{doc.filename}》。\n"
        f"下载地址：{doc.download_url}\n"
        f"请用中文告知用户文件已就绪，并附上该下载地址（原样保留路径，方便前端识别）；"
        f"用户可在对话里点「预览」或「下载」。"
    )


def _esc(text: str) -> str:
    import html

    return html.escape(text or "", quote=True)


def _preview_docx_html(path: Path) -> str:
    from docx import Document

    document = Document(str(path))
    chunks: list[str] = ['<div class="doc-preview docx">']
    for p in document.paragraphs:
        text = (p.text or "").strip()
        if not text:
            continue
        style = (p.style.name or "") if p.style else ""
        if "Heading 1" in style or style == "Title":
            chunks.append(f"<h1>{_esc(text)}</h1>")
        elif "Heading 2" in style:
            chunks.append(f"<h2>{_esc(text)}</h2>")
        elif "Heading 3" in style:
            chunks.append(f"<h3>{_esc(text)}</h3>")
        elif text.startswith("• ") or text.startswith("- "):
            chunks.append(f"<li>{_esc(text[2:].strip())}</li>")
        else:
            chunks.append(f"<p>{_esc(text)}</p>")
    for table in document.tables:
        chunks.append("<table>")
        for row in table.rows:
            chunks.append("<tr>")
            for cell in row.cells:
                chunks.append(f"<td>{_esc((cell.text or '').strip())}</td>")
            chunks.append("</tr>")
        chunks.append("</table>")
    chunks.append("</div>")
    return "\n".join(chunks)


def _preview_pptx_html(path: Path) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    chunks: list[str] = ['<div class="doc-preview pptx">']
    for i, slide in enumerate(prs.slides, start=1):
        chunks.append(f'<section class="slide"><h2>幻灯片 {i}</h2>')
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip() or (para.text or "").strip()
                    if text:
                        chunks.append(f"<p>{_esc(text)}</p>")
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                chunks.append("<table>")
                for row in shape.table.rows:
                    chunks.append("<tr>")
                    for cell in row.cells:
                        chunks.append(f"<td>{_esc((cell.text or '').strip())}</td>")
                    chunks.append("</tr>")
                chunks.append("</table>")
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                note = (slide.notes_slide.notes_text_frame.text or "").strip()
                if note:
                    chunks.append(f'<p class="notes"><em>备注：{_esc(note)}</em></p>')
        except Exception:  # noqa: BLE001
            pass
        chunks.append("</section>")
    chunks.append("</div>")
    return "\n".join(chunks)


def _preview_md_html(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    chunks: list[str] = ['<div class="doc-preview md">']
    for line in raw.splitlines():
        s = line.rstrip()
        if not s:
            chunks.append("<br/>")
            continue
        if s.startswith("### "):
            chunks.append(f"<h3>{_esc(s[4:])}</h3>")
        elif s.startswith("## "):
            chunks.append(f"<h2>{_esc(s[3:])}</h2>")
        elif s.startswith("# "):
            chunks.append(f"<h1>{_esc(s[2:])}</h1>")
        elif s.startswith("- ") or s.startswith("* "):
            chunks.append(f"<li>{_esc(s[2:])}</li>")
        elif s.startswith("---"):
            chunks.append("<hr/>")
        else:
            chunks.append(f"<p>{_esc(s)}</p>")
    chunks.append("</div>")
    return "\n".join(chunks)


def build_generated_preview(
    file_id: str,
    *,
    username: str,
    elevated: bool = False,
) -> dict:
    """Return {file_id, filename, kind, html} for in-app preview."""
    path, filename, _media = resolve_generated_file(
        file_id, username=username, elevated=elevated
    )
    meta = load_generated_meta(file_id) or {}
    kind = str(meta.get("kind") or path.suffix.lstrip(".") or "docx")
    if kind == "docx":
        html = _preview_docx_html(path)
    elif kind == "pptx":
        html = _preview_pptx_html(path)
    elif kind == "md":
        html = _preview_md_html(path)
    else:
        raise ValueError(f"暂不支持预览该类型：{kind}")
    return {
        "file_id": file_id,
        "filename": filename,
        "kind": kind,
        "html": html,
    }
