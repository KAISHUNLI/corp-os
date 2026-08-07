"""Extract searchable text from uploaded office/image files."""

from __future__ import annotations

from pathlib import Path


# Enterprise allow-list: text-extractable knowledge / expense / template materials.
TEXT_EXTENSIONS = {".txt", ".md"}
TABULAR_EXTENSIONS = {".csv", ".xlsx"}
PDF_EXTENSIONS = {".pdf"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}
DOC_EXTENSIONS = {".docx"}
PPT_EXTENSIONS = {".pptx"}

SUPPORTED_EXTENSIONS = (
    TEXT_EXTENSIONS
    | TABULAR_EXTENSIONS
    | PDF_EXTENSIONS
    | IMAGE_EXTENSIONS
    | DOC_EXTENSIONS
    | PPT_EXTENSIONS
)


def extract_text(path: Path, *, override: str | None = None) -> str:
    if override and override.strip():
        return override.strip()

    ext = path.suffix.lower()
    if ext in TEXT_EXTENSIONS:
        return path.read_text(encoding="utf-8", errors="ignore").strip()
    if ext == ".csv":
        return path.read_text(encoding="utf-8", errors="ignore").strip()
    if ext == ".xlsx":
        return _extract_xlsx(path)
    if ext in PDF_EXTENSIONS:
        return _extract_pdf(path)
    if ext in IMAGE_EXTENSIONS:
        return _extract_image(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    return ""


def _extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("缺少 openpyxl，无法解析 Excel。请 pip install openpyxl") from exc

    wb = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"# 工作表：{sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append("\t".join(cells))
    wb.close()
    return "\n".join(parts).strip()


def _extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("缺少 pypdf，无法解析 PDF。请 pip install pypdf") from exc

    reader = PdfReader(str(path))
    parts: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            parts.append(f"# 第 {i} 页\n{text}")
    return "\n\n".join(parts).strip()


def _extract_image(path: Path) -> str:
    # Optional OCR: works when pytesseract + system tesseract are available.
    try:
        from PIL import Image
        import pytesseract

        image = Image.open(path)
        text = pytesseract.image_to_string(image, lang="chi_sim+eng")
        text = (text or "").strip()
        if text:
            return text
    except Exception:
        pass

    return (
        f"[已上传图片：{path.name}。未识别到文字（可安装 tesseract OCR，"
        "或在上传备注/文本框中粘贴图片上的关键内容以便检索与报销预审）。]"
    )


def _extract_docx(path: Path) -> str:
    try:
        import docx
    except ImportError:
        return f"[已上传 Word：{path.name}。未安装 python-docx，未能抽取正文。]"

    document = docx.Document(str(path))
    parts = [p.text.strip() for p in document.paragraphs if p.text and p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append("\t".join(cells))
    text = "\n".join(parts).strip()
    if text:
        return text
    return f"[已上传 Word：{path.name}，未解析到正文。可在备注中补充关键内容。]"


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return f"[已上传 PPT：{path.name}。未安装 python-pptx，未能抽取正文。]"

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip() or (para.text or "").strip()
                    if text:
                        slide_lines.append(text)
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                    if cells:
                        slide_lines.append("\t".join(cells))
        # notes
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                note = (slide.notes_slide.notes_text_frame.text or "").strip()
                if note:
                    slide_lines.append(f"（备注）{note}")
        except Exception:  # noqa: BLE001
            pass
        if slide_lines:
            parts.append(f"# 幻灯片：{i}\n" + "\n".join(slide_lines))

    text = "\n\n".join(parts).strip()
    if text:
        return text
    return f"[已上传 PPT：{path.name}，未解析到正文。可在备注中补充关键内容。]"
